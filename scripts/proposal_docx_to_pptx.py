#!/usr/bin/env python3
"""
Build a short, presentation-friendly .pptx from the research proposal .docx.

Reads structure (headings, paragraphs, tables) from the Word file and maps each
major section to one or more slides with compressed bullet points.

Requires: python-docx, python-pptx
  pip install python-docx python-pptx

Examples:
  python scripts/proposal_docx_to_pptx.py
  python scripts/proposal_docx_to_pptx.py -i "C:\\path\\proposal.docx" -o deck.pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path
from typing import Iterator


def ensure_packages() -> None:
    for mod, pip_name in (("docx", "python-docx"), ("pptx", "python-pptx")):
        try:
            __import__(mod)
        except ImportError:
            print(f"Installing {pip_name} …", file=sys.stderr)
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", pip_name],
                stdout=sys.stderr,
            )
            __import__(mod)


def paragraph_heading_level(paragraph) -> int | None:
    """Return 0 for Title, 1–9 for Heading n, None if body text."""
    name = (paragraph.style.name or "").strip()
    if name == "Title":
        return 0
    if name.startswith("Heading"):
        parts = name.split()
        if len(parts) >= 2 and parts[-1].isdigit():
            return int(parts[-1])
    return None


def iter_block_items(document):
    """Yield Paragraph and Table objects in document order (body)."""
    from docx.document import Document as DocType
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def strip_citations(text: str) -> str:
    return re.sub(r"\s*\[\d+\](?:\s*,\s*\[\d+\])*\s*", " ", text)


def compress_line(text: str, max_len: int = 130) -> str:
    text = strip_citations(text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return ""
    # Unify fancy quotes
    text = text.replace("“", '"').replace("”", '"').replace("’", "'")
    if len(text) <= max_len:
        return text
    cut = text[: max_len - 1]
    sp = cut.rfind(" ")
    if sp > 40:
        cut = cut[:sp]
    return cut.rstrip(",;:") + "…"


def paragraph_to_bullets(text: str, max_bullets: int = 4) -> list[str]:
    """Turn a paragraph into a few short bullets (by sentence)."""
    text = text.strip()
    if not text:
        return []
    # Split on sentence boundaries (simple)
    parts = re.split(r"(?<=[.!?])\s+", text)
    bullets: list[str] = []
    for p in parts:
        p = p.strip()
        if not p:
            continue
        line = compress_line(p)
        if line and line not in bullets:
            bullets.append(line)
        if len(bullets) >= max_bullets:
            break
    if not bullets and text:
        bullets.append(compress_line(text))
    return bullets


def table_to_bullets(table, max_rows: int = 8) -> list[str]:
    rows: list[str] = []
    for i, row in enumerate(table.rows):
        if i >= max_rows:
            rows.append("… (table continues in document)")
            break
        cells = [c.text.strip() for c in row.cells if c.text.strip()]
        if not cells:
            continue
        rows.append(" · ".join(cells))
    return rows


def parse_docx(path: Path) -> tuple[str, str, list[tuple[str, list[str]]]]:
    """
    Returns (main_title, subtitle, sections) where each section is (heading, bullets).
    """
    ensure_packages()
    from docx import Document
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph

    doc = Document(str(path))
    main_title = "Research proposal"
    subtitle = ""
    sections: list[tuple[str, list[str]]] = []
    current_heading: str | None = None
    current_bullets: list[str] = []

    def flush():
        nonlocal current_heading, current_bullets
        if current_heading is None:
            current_bullets.clear()
            return
        # Drop empty sections except keep structure for tables-only sections
        cleaned = [b for b in current_bullets if b]
        if cleaned or current_heading.lower().startswith(("11.", "14.")):
            sections.append((current_heading, cleaned))
        current_heading = None
        current_bullets = []

    for block in iter_block_items(doc):
        if isinstance(block, DocxParagraph):
            p = block
            text = p.text.strip()
            if not text:
                continue
            level = paragraph_heading_level(p)
            if level == 0:
                main_title = text
                continue
            if level == 1:
                flush()
                current_heading = text
                current_bullets = []
                continue
            if level == 2 and current_heading:
                current_bullets.append(compress_line(text, max_len=100))
                continue
            # Body or list: first line after title page "Research Proposal"
            if main_title and not sections and not current_heading and text == "Research Proposal":
                subtitle = text
                continue
            if current_heading:
                for b in paragraph_to_bullets(text, max_bullets=3):
                    if b not in current_bullets:
                        current_bullets.append(b)
            continue

        if isinstance(block, DocxTable) and current_heading:
            for b in table_to_bullets(block):
                if b not in current_bullets:
                    current_bullets.append(b)

    flush()

    # Merge duplicate title section "1. Title" body into cleaner first content slide
    return main_title, subtitle, sections


def split_bullets_for_slides(bullets: list[str], per_slide: int = 6) -> list[list[str]]:
    if not bullets:
        return [[]]
    return [bullets[i : i + per_slide] for i in range(0, len(bullets), per_slide)]


def build_presentation(
    main_title: str,
    subtitle: str,
    sections: list[tuple[str, list[str]]],
    out_path: Path,
    bullets_per_slide: int = 6,
) -> None:
    ensure_packages()
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title slide (blank layout: manual title box) ---
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = compress_line(main_title, max_len=200)
    p0.font.size = Pt(32)
    p0.font.bold = True
    if subtitle:
        p1 = tf.add_paragraph()
        p1.text = subtitle
        p1.font.size = Pt(18)
        p1.space_before = Pt(12)
    sub2 = tf.add_paragraph()
    sub2.text = "Short overview (from Word proposal)"
    sub2.font.size = Pt(14)
    sub2.font.italic = True
    sub2.space_before = Pt(18)

    title_and_content = prs.slide_layouts[1]

    for heading, bullets in sections:
        # Special short slides
        if heading.startswith("15.") and bullets:
            bullets = [
                "IEEE-style sources (see full list in Word)",
                compress_line(bullets[0], 90),
                compress_line(bullets[1], 90) if len(bullets) > 1 else "",
                compress_line(bullets[2], 90) if len(bullets) > 2 else "",
                "+ additional references in document",
            ]
            bullets = [b for b in bullets if b]

        slide_title = compress_line(heading.replace("14. Budget", "14. Budget (indicative)"), max_len=70)
        chunks = split_bullets_for_slides(bullets, bullets_per_slide)
        for idx, chunk in enumerate(chunks):
            sl = prs.slides.add_slide(title_and_content)
            shapes = sl.shapes
            t = shapes.title
            t.text = f"{slide_title} ({idx + 1}/{len(chunks)})" if len(chunks) > 1 else slide_title
            body = shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True
            use = chunk if chunk else ["(See full proposal.)"]
            tf.text = use[0]
            tf.paragraphs[0].font.size = Pt(16 if len(use) <= 5 else 14)
            tf.paragraphs[0].space_after = Pt(4)
            for line in use[1:]:
                para = tf.add_paragraph()
                para.text = line
                para.font.size = Pt(16 if len(use) <= 5 else 14)
                para.level = 0
                para.space_after = Pt(4)

    # Closing slide
    sl = prs.slides.add_slide(title_and_content)
    sl.shapes.title.text = "Takeaways"
    tf = sl.shapes.placeholders[1].text_frame
    tf.clear()
    closing = [
        "Rural-linked context → slower adoption; train + hybrid (voice) options.",
        "Main gaps today: no digital dispatch, manual calls, unclear fares, wasted time/fuel.",
        "Stakeholders see value if fairness, safety, and low commissions are credible.",
        "Next: finish fieldwork, then pilot design if feasibility looks positive.",
    ]
    tf.text = closing[0]
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].space_after = Pt(6)
    for line in closing[1:]:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(18)
        para.space_after = Pt(6)

    prs.save(str(out_path))


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    default_in = root / "Feasibility_Study_Ride_Hailing_Debremarkos_Proposal.docx"
    default_out = root / "Feasibility_Study_Ride_Hailing_Debremarkos_Overview.pptx"

    ap = argparse.ArgumentParser(description="Convert proposal .docx to short .pptx")
    ap.add_argument("-i", "--input", type=Path, default=default_in, help="Input .docx path")
    ap.add_argument("-o", "--output", type=Path, default=default_out, help="Output .pptx path")
    ap.add_argument("--bullets-per-slide", type=int, default=6, help="Max bullets before splitting slide")
    args = ap.parse_args()

    in_path: Path = args.input.expanduser().resolve()
    if not in_path.is_file():
        print(f"Input file not found: {in_path}", file=sys.stderr)
        print("Generate it first with: python scripts/generate_ride_hailing_research_proposal.py", file=sys.stderr)
        sys.exit(1)

    main_title, subtitle, sections = parse_docx(in_path)
    out_path: Path = args.output.expanduser().resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    build_presentation(main_title, subtitle, sections, out_path, args.bullets_per_slide)
    print(f"Saved: {out_path} ({len(sections)} sections -> slides)")


if __name__ == "__main__":
    main()
